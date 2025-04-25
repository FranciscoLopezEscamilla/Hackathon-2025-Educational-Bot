from pptx import Presentation
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.util import Inches, Pt
import os

TEXT_BLOCK_HEIGHT = Inches(0.8)
IMAGE_HEIGHT = Inches(2)
IMAGE_SPACING = Inches(0.2)
MAX_CONTENT_HEIGHT = Inches(7.0)
LOCAL_IMG_DIR="../assets/generated_images/2025-04-21-11-57-47"

class PptGenerator:

    def generate_ppt(doc_data, output_path: str):
        prs = Presentation()

        for slide_data in doc_data.pages:
            slide = prs.slides.add_slide(prs.slide_layouts[5])

            # Customizing text dimensions
            num_blocks = len(slide_data.text_items or [])
            text_height_estimate = num_blocks * TEXT_BLOCK_HEIGHT
            textbox = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(8.5), text_height_estimate)
            tf = textbox.text_frame
            tf.word_wrap = True
            tf.auto_size = None#MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

            for i, element in enumerate(slide_data.text_items or []):
                para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                para.text = element.content
                para.space_after = Pt(10)
                font = para.font

                if element.type == "title":
                    font.size = Pt(32)
                    font.bold = True
                elif element.type == "subtitle":
                    font.size = Pt(24)
                    font.italic = True
                elif element.type == "paragraph":
                    font.size = Pt(16)
                else:
                    font.size = Pt(14)

            # Customizing images
            image_top = Inches(0.5) + text_height_estimate + Inches(0.5)

            for img_name in slide_data.images or []:
                image_path = os.path.join(LOCAL_IMG_DIR, img_name)
                if os.path.exists(image_path):
                    # Check if this image fits on the current slide
                    if image_top + IMAGE_HEIGHT > MAX_CONTENT_HEIGHT:
                        # Create a new slide and reset
                        slide = prs.slides.add_slide(prs.slide_layouts[5])
                        image_top = Inches(0.5)

                    try:
                        slide.shapes.add_picture(image_path, Inches(0.5), image_top, height=IMAGE_HEIGHT)
                        image_top += IMAGE_HEIGHT + IMAGE_SPACING
                    except Exception as e:
                        print(f"❌ Error adding new image: {e}")

        prs.save(output_path)

        